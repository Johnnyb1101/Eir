from pptx import Presentation
from pptx.util import Inches

def render_deck(deck, path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for s in deck.slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        slide.shapes.title.text = s.title
        slide.shapes.title.left = Inches(0.5)
        slide.shapes.title.top = Inches(0.3)
        slide.shapes.title.width = Inches(12.3)
        slide.shapes.title.height = Inches(1.2)

        ph = slide.placeholders[1]
        ph.left = Inches(0.5)
        ph.top = Inches(1.6)
        ph.width = Inches(12.3)
        ph.height = Inches(5.6)

        body = ph.text_frame
        body.text = s.bullets[0]
        for b in s.bullets[1:]:
            body.add_paragraph().text = b
        body.fit_text()

        notes = slide.notes_slide.notes_text_frame
        notes.text = s.speaker_notes
        notes.add_paragraph().text = ""
        for c in s.sources:
            notes.add_paragraph().text = f"{c.source}, {c.section}, pp. {c.pages}"

    refs = []
    for s in deck.slides:
        for c in s.sources:
            if c.source not in refs:
                refs.append(c.source)

    if refs:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "References"
        slide.shapes.title.left = Inches(0.5)
        slide.shapes.title.top = Inches(0.3)
        slide.shapes.title.width = Inches(12.3)
        slide.shapes.title.height = Inches(1.2)

        ph = slide.placeholders[1]
        ph.left = Inches(0.5)
        ph.top = Inches(1.6)
        ph.width = Inches(12.3)
        ph.height = Inches(5.6)

        body = ph.text_frame
        body.text = refs[0]
        for r in refs[1:]:
            body.add_paragraph().text = r
        body.fit_text()

    prs.save(path)